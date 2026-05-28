import warnings
# Suppress EOL and SSL version warnings cleanly at the compiler layer
warnings.filterwarnings("ignore", category=FutureWarning)
warnings.filterwarnings("ignore", message=".*OpenSSL.*")

import streamlit as st
from google import genai
from google.genai import types
from PIL import Image
import PyPDF2
from pptx import Presentation
import io
import hashlib
import json
from pymongo import MongoClient

# 1. Setup Page Configuration
st.set_page_config(page_title="VFSTR Academic Command Center", page_icon="🎓", layout="wide")

# 2. Initialize Gemini Client with your API Key
try:
    client = genai.Client(api_key="AIzaSyCkFJRM0zGlo-wVBZBIhM8F5c0Ac3t2AqA")
except Exception as e:
    st.error(f"Failed to initialize Gemini Client: {e}")
    st.stop()

# ==========================================
# MONGODB LOCAL DATABASE CONNECTION
# ==========================================
MONGO_URI = "mongodb://localhost:27017/"

@st.cache_resource
def init_mongodb():
    try:
        db_client = MongoClient(MONGO_URI, serverSelectionTimeoutMS=5000)
        db = db_client["vfstr_academic_hub"]
        return db["users"]
    except Exception as error:
        st.sidebar.error(f"Database offline: {error}")
        return None

users_collection = init_mongodb()

def hash_password(password):
    return hashlib.sha256(str.encode(password)).hexdigest()

# Dynamic DB Helper Functions for Real-Time Synchronization
def sync_telemetry_to_db():
    """Pushes temporary session updates out to the persistent MongoDB collection."""
    if st.session_state.authenticated and users_collection is not None:
        users_collection.update_one(
            {"username": st.session_state.username},
            {"$set": {
                "topics_planned": st.session_state.topics_planned,
                "quizzes_generated": st.session_state.quizzes_generated,
                "codes_analyzed": st.session_state.codes_analyzed,
                "last_analyzed_topic": st.session_state.last_analyzed_topic
            }}
        )

# Initialize Session States for Auth & Local Telemetry
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False
if "username" not in st.session_state:
    st.session_state.username = ""
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "quizzes_generated" not in st.session_state:
    st.session_state.quizzes_generated = 0
if "codes_analyzed" not in st.session_state:
    st.session_state.codes_analyzed = 0
if "topics_planned" not in st.session_state:
    st.session_state.topics_planned = 0
if "last_analyzed_topic" not in st.session_state:
    st.session_state.last_analyzed_topic = ""

# Persistent Quiz Storage across form submissions
if "current_quiz" not in st.session_state:
    st.session_state.current_quiz = None
if "quiz_submitted" not in st.session_state:
    st.session_state.quiz_submitted = False

# Document Parsers
def extract_pdf_text(file_bytes):
    pdf_file = io.BytesIO(file_bytes)
    reader = PyPDF2.PdfReader(pdf_file)
    return "".join([page.extract_text() or "" for page in reader.pages])

def extract_pptx_text(file_bytes):
    pptx_file = io.BytesIO(file_bytes)
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# ==========================================
# AUTHENTICATION INTERFACE (MONGODB BACKED)
# ==========================================
if not st.session_state.authenticated:
    st.title("🦅 VFSTR AI Hub: Academic Command Center")
    st.write("Welcome! Please Register .")
    
    auth_mode = st.radio("Choose Action:", ["Sign Up / Register", "Login"])
    
    if users_collection is None:
        st.error("⚠️ Unable to link up database pipeline. Please make sure your local MongoDB instance is running.")
        st.stop()
        
    if auth_mode == "Sign Up / Register":
        st.subheader("📝 Live Student Registration")
        new_username = st.text_input("Create Username / Regd No:")
        new_password = st.text_input("Create Password:", type="password")
        confirm_password = st.text_input("Confirm Password:", type="password")
        
        if st.button("Register Account 🎉"):
            if not new_username.strip() or not new_password.strip():
                st.warning("Input channels cannot sit empty.")
            elif new_password != confirm_password:
                st.error("Passwords do not match.")
            else:
                existing_user = users_collection.find_one({"username": new_username})
                if existing_user:
                    st.error("User ID already logged in system directories.")
                else:
                    user_payload = {
                        "username": new_username,
                        "password": hash_password(new_password),
                        "topics_planned": 0,
                        "quizzes_generated": 0,
                        "codes_analyzed": 0,
                        "last_analyzed_topic": ""
                    }
                    users_collection.insert_one(user_payload)
                    st.success("Registration completely written to MongoDB! Please switch the action toggle above to 'Login' to access your dashboard.")
                    
    elif auth_mode == "Login":
        st.subheader("🔑 Student Cloud Login")
        username_input = st.text_input("Regd No / Username:")
        password_input = st.text_input("Password:", type="password")
        
        if st.button("Authenticate Identity 🚀"):
            hashed_p = hash_password(password_input)
            user_record = users_collection.find_one({"username": username_input})
            
            if user_record and user_record["password"] == hashed_p:
                st.session_state.authenticated = True
                st.session_state.username = username_input
                
                # Load saved metrics directly back into active session values
                st.session_state.topics_planned = user_record.get("topics_planned", 0)
                st.session_state.quizzes_generated = user_record.get("quizzes_generated", 0)
                st.session_state.codes_analyzed = user_record.get("codes_analyzed", 0)
                st.session_state.last_analyzed_topic = user_record.get("last_analyzed_topic", "")
                
                st.success(f"Welcome back, {username_input}!")
                st.rerun()
            else:
                st.error("Invalid Username or Password credential pairs.")

# ==========================================
# AUTHENTICATED COMMAND UTILITIES
# ==========================================
else:
    st.sidebar.title("🦅 VFSTR AI Hub")
    st.sidebar.markdown(f"**👤 Connected Account:** `{st.session_state.username}`")
    
    if st.sidebar.button("🔒 Logout"):
        st.session_state.authenticated = False
        st.session_state.username = ""
        st.session_state.chat_history = []
        st.session_state.quizzes_generated = 0
        st.session_state.codes_analyzed = 0
        st.session_state.topics_planned = 0
        st.session_state.last_analyzed_topic = ""
        st.session_state.current_quiz = None
        st.session_state.quiz_submitted = False
        st.rerun()
        
    st.sidebar.markdown("---")
    choice = st.sidebar.radio("Navigate Tools:", [
        "💬 Interactive AI Study Helper", 
        "📅 Smart Syllabus Planner", 
        "📝 Automated Quiz Generator",
        "💻 Advanced Code Debugger",
        "📊 Live Student Analytics"
    ])

    # FEATURE 1: Chat Helper
    if choice == "💬 Interactive AI Study Helper":
        st.title("💬 Interactive AI Academic Study Helper")
        st.write(f"Hello `{st.session_state.username}`, ask Meghanadh anything regarding your subjects or career paths.")
        
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["text"])
                
        if user_query := st.chat_input("Ask Meghanadh anything..."):
            with st.chat_message("user"):
                st.markdown(user_query)
            st.session_state.chat_history.append({"role": "user", "text": user_query})
            
            with st.chat_message("assistant"):
                with st.spinner("Meghanadh is thinking..."):
                    try:
                        response = client.models.generate_content(
                            model='gemini-2.5-flash',
                            contents=user_query,
                            config=types.GenerateContentConfig(
                                system_instruction="You are Meghanadh, a brilliant, helpful, and sharp Computer Science mentor at Vignan's University (VFSTR). Guide the student step-by-step with clear technical insights and an encouraging tone.",
                            )
                        )
                        st.markdown(response.text)
                        st.session_state.chat_history.append({"role": "assistant", "text": response.text})
                    except Exception as e:
                        st.error(f"Chat Error: {e}")

    # FEATURE 2: Syllabus Planner
    elif choice == "📅 Smart Syllabus Planner":
        st.title("📅 Smart Syllabus & Note Scheduler")
        st.write("Upload or paste your course topics to generate an optimized, step-by-step study plan.")
        source_type = str(st.radio("Input Method:", ["Paste Text", "Upload File (.txt)"]))
        syllabus_text = ""
        if source_type == "Paste Text":
            syllabus_text = st.text_area("Paste your syllabus or lecture topics here:", height=150)
        else:
            uploaded_file = st.file_uploader("Choose a text file", type=["txt"])
            if uploaded_file is not None:
                syllabus_text = uploaded_file.read().decode("utf-8")
        days = st.slider("Over how many days do you want to study this?", 1, 30, 7)
        if st.button("Generate Study Plan 🚀"):
            if not syllabus_text.strip():
                st.warning("Please provide some syllabus text first!")
            else:
                with st.spinner("Gemini is analyzing your material..."):
                    prompt = f"Analyze this material and break it down into an efficient {days}-day study schedule. Provide Focus Topic, Core Concepts, and a Quick Practice Task for each day.\n\nMaterial:\n{syllabus_text}"
                    try:
                        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                        st.success("Plan ready!")
                        st.markdown(response.text)
                        st.session_state.topics_planned += 1
                        st.session_state.last_analyzed_topic = syllabus_text[:50] + "..."
                        sync_telemetry_to_db()
                    except Exception as e:
                        st.error(f"Error: {e}")

    # FEATURE 3: Quiz Generator
    elif choice == "📝 Automated Quiz Generator":
        st.title("📝 Automated Quiz & Assessment Generator")
        st.write("Upload documents or paste transcripts to construct a **one-attempt examination module**.")
        
        input_mode = st.radio("Source Material Format:", ["Direct Text Input", "Upload Reference Document"])
        quiz_material = ""
        uploaded_image = None
        num_questions = st.slider("Number of questions to generate", 3, 10, 5)
        
        if input_mode == "Direct Text Input":
            quiz_material = st.text_area("Paste study transcripts here:", height=150)
        else:
            uploaded_doc = st.file_uploader("Choose a document module", type=["pdf", "pptx", "png", "jpg", "jpeg"])
            if uploaded_doc is not None:
                file_type = uploaded_doc.name.split(".")[-1].lower()
                with st.spinner("Extracting contents from source file..."):
                    if file_type in ["png", "jpg", "jpeg"]:
                        uploaded_image = Image.open(uploaded_doc)
                        st.image(uploaded_image, caption="Uploaded Image Reference Structure", width=400)
                    elif file_type == "pdf":
                        quiz_material = extract_pdf_text(uploaded_doc.read())
                        st.success(f"Successfully extracted text from {uploaded_doc.name}!")
                    elif file_type == "pptx":
                        quiz_material = extract_pptx_text(uploaded_doc.read())
                        st.success(f"Successfully extracted slide data from {uploaded_doc.name}!")

        if st.button("Generate Dynamic Quiz 🧠"):
            json_schema = {
                "type": "OBJECT",
                "properties": {
                    "questions": {
                        "type": "ARRAY",
                        "items": {
                            "type": "OBJECT",
                            "properties": {
                                "question_text": {"type": "STRING"},
                                "options": {"type": "ARRAY", "items": {"type": "STRING"}},
                                "correct_answer": {"type": "STRING"},
                                "explanation": {"type": "STRING"}
                            },
                            "required": ["question_text", "options", "correct_answer", "explanation"]
                        }
                    }
                },
                "required": ["questions"]
            }
            
            base_prompt = f"Based on the source material context, compile exactly {num_questions} unique multiple choice questions. The 'correct_answer' string must exactly match one of the items written inside the 'options' list array."
            
            try:
                with st.spinner("Gemini is constructing structured evaluation variables..."):
                    if uploaded_image:
                        response = client.models.generate_content(
                            model='gemini-2.5-flash',
                            contents=[uploaded_image, base_prompt],
                            config=types.GenerateContentConfig(response_mime_type="application/json", response_schema=json_schema)
                        )
                        st.session_state.last_analyzed_topic = f"Visual Quiz ({uploaded_doc.name})"
                    elif quiz_material.strip():
                        response = client.models.generate_content(
                            model='gemini-2.5-flash',
                            contents=f"{base_prompt}\n\nReference Material:\n{quiz_material}",
                            config=types.GenerateContentConfig(response_mime_type="application/json", response_schema=json_schema)
                        )
                        st.session_state.last_analyzed_topic = f"Document Quiz ({uploaded_doc.name if 'uploaded_doc' in locals() else 'Raw Text'})"
                    else:
                        st.warning("Please provide context text or files first.")
                        response = None

                    if response:
                        st.session_state.current_quiz = json.loads(response.text)["questions"]
                        st.session_state.quizzes_generated += 1
                        st.session_state.quiz_submitted = False
                        sync_telemetry_to_db()
                        st.rerun()
            except Exception as err:
                st.error(f"Failed to generate structured quiz matrix: {err}")

        if st.session_state.current_quiz:
            st.markdown("---")
            st.subheader("✍️ Live Test Session")
            
            user_responses = {}
            for i, q in enumerate(st.session_state.current_quiz):
                st.markdown(f"**Q{i+1}: {q['question_text']}**")
                display_options = ["-- Select your Answer --"] + q["options"]
                
                user_responses[i] = st.radio(
                    f"Options for Q{i+1}:", 
                    options=display_options, 
                    key=f"q_{i}_input",
                    label_visibility="collapsed",
                    disabled=st.session_state.quiz_submitted
                )
                
                if st.session_state.quiz_submitted:
                    if user_responses[i] == q["correct_answer"]:
                        st.success(f"✅ **Correct!** {q['explanation']}")
                    else:
                        st.error(f"❌ **Incorrect.** Correct Answer: {q['correct_answer']} \n\n {q['explanation']}")
                st.markdown("<br>", unsafe_allow_html=True)
            
            if not st.session_state.quiz_submitted:
                if st.button("🔒 Submit Test & View Results"):
                    unanswered = [idx for idx, ans in user_responses.items() if ans == "-- Select your Answer --"]
                    if unanswered:
                        st.error(f"Please answer all questions before submitting! (Unanswered: Question(s) {[idx+1 for idx in unanswered]})")
                    else:
                        st.session_state.quiz_submitted = True
                        st.rerun()
            else:
                total_q = len(st.session_state.current_quiz)
                final_correct = sum(1 for idx, ans in user_responses.items() if ans == st.session_state.current_quiz[idx]["correct_answer"])
                st.success(f"📊 **Test Complete! Final Score: {final_correct} / {total_q}**")
                if st.button("🔄 Take New Quiz"):
                    st.session_state.current_quiz = None
                    st.session_state.quiz_submitted = False
                    st.rerun()

    # FEATURE 4: Advanced Code Debugger
    elif choice == "💻 Advanced Code Debugger":
        st.title("💻 Advanced Code Explainer & Debugger")
        lang = st.selectbox("Select Language:", ["Python", "Java", "C++", "JavaScript", "HTML/CSS"])
        code_input = st.text_area("Paste code here:", height=200)
        purpose = st.radio("Action:", ["Debug & Fix Code", "Explain How Code Works"])
        if st.button("Process Code"):
            if not code_input.strip():
                st.warning("Please paste some code first!")
            else:
                with st.spinner("Analyzing structure..."):
                    prompt = f"Analyze this {lang} code. Act according to request: {purpose}.\n\nCode:\n{code_input}"
                    try:
                        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                        st.markdown(response.text)
                        st.session_state.codes_analyzed += 1
                        sync_telemetry_to_db()
                    except Exception as e:
                        st.error(f"Error: {e}")

    # FEATURE 5: Live Analytics Dashboard
    elif choice == "📊 Live Student Analytics":
        st.title("📊 Live Personal AI Insights Dashboard")
        st.write("This dashboard monitors your live session interactions and leverages Gemini to synthesize a real-time study strategy review.")
        st.markdown("---")
        
        if users_collection is not None:
            fresh_record = users_collection.find_one({"username": st.session_state.username})
            if fresh_record:
                db_topics = fresh_record.get("topics_planned", 0)
                db_quizzes = fresh_record.get("quizzes_generated", 0)
                db_codes = fresh_record.get("codes_analyzed", 0)
                db_last_topic = fresh_record.get("last_analyzed_topic", "")
            else:
                db_topics, db_quizzes, db_codes, db_last_topic = 0, 0, 0, ""
        else:
            db_topics, db_quizzes, db_codes, db_last_topic = 0, 0, 0, ""

        # Symmetrical native horizontal layout system
        col1, col2, col3 = st.columns(3, gap="large")
        with col1:
            st.markdown("### 📋")
            st.metric(label="Syllabus Units Planned", value=f"{db_topics} Units")
        with col2:
            st.markdown("### 📝")
            st.metric(label="Quizzes Generated", value=f"{db_quizzes} Sessions")
        with col3:
            st.markdown("### 💻")
            st.metric(label="Code Debugs Run", value=f"{db_codes} Scripts")
        
        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown("### 📈 Live AI Strategy Assessment")
        
        total_interactions = db_topics + db_quizzes + db_codes
        if total_interactions == 0:
            st.warning("⚠️ **Telemetry stream is currently empty.**")
            st.write(f"Hello `{st.session_state.username}`, the AI engine requires session telemetry data to construct a performance report. Please navigate to the tools on the left to start your study tracker.")
            st.markdown("#### 🎯 Session Activation Milestones")
            st.progress(0)
            st.caption("0 / 3 Key utilities initialized during this run.")
        else:
            with st.spinner("Gemini is compiling your strategy report metrics..."):
                last_topic = db_last_topic if db_last_topic else "No external files loaded yet"
                
                prompt = f"Act as an elite university academic advisor analyzing real-time student activity. Provide a personalized study strategy review based on these metrics:\n- Syllabus Modules Planned: {db_topics}\n- Practice Quizzes Generated: {db_quizzes}\n- Code Interventions: {db_codes}\n- Last Active Material: {last_topic}\n\nStructure your output with headers for 'Current Session Strengths' and 'Recommended Optimization'. Do not repeat exact phrasing from previous runs."
                try:
                    response = client.models.generate_content(
                        model='gemini-2.5-flash', 
                        contents=prompt,
                        config=types.GenerateContentConfig(
                            temperature=0.85, # Injects deep dynamic variety on text execution sweeps
                            top_p=0.95
                        )
                    )
                    st.success("✨ Live AI Assessment Generated Instantly via Gemini 2.5 Flash")
                    st.info(response.text)
                except Exception as e:
                    st.info("💡 **Strategy Report Server Refreshing:** The system analytics engine is updating the token buffers. Press 'R' or toggle back in a moment to populate your strategy card!")